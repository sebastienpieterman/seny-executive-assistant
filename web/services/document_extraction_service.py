"""
DocumentExtractionService — Phase 28 File Upload & Document Intelligence

Stateless/infrastructure service (no user_id) that detects file type and
extracts text content from uploaded files. Supports PDF, Word, PowerPoint,
CSV, plain text, Markdown, HTML, and images.
"""
import base64
import csv
import io
import logging
import os
import re

logger = logging.getLogger(__name__)


class DocumentExtractionService:
    """
    Stateless service for detecting file types and extracting text content.

    No user_id in constructor — this is infrastructure-level, user-agnostic.
    All extraction errors are caught and returned as descriptive error text
    rather than raised exceptions.
    """

    # Supported image types for multimodal Claude passthrough
    IMAGE_EXTENSIONS = {".png", ".jpg", ".jpeg", ".webp", ".gif"}

    # Map extension -> logical file type
    EXTENSION_MAP = {
        ".pdf": "pdf",
        ".docx": "word",
        ".pptx": "powerpoint",
        ".ppt": "powerpoint",
        ".csv": "csv",
        ".txt": "text",
        ".md": "text",
        ".html": "html",
        ".htm": "html",
    }

    async def extract(self, filename: str, content_bytes: bytes, content_type: str) -> dict:
        """
        Detect file type and extract content from the given bytes.

        Args:
            filename: Original filename (used for extension-based type detection).
            content_bytes: Raw file bytes.
            content_type: MIME type reported by the client (fallback if no extension).

        Returns:
            dict with keys:
                file_name: str — original filename
                file_type: str — "pdf", "word", "powerpoint", "csv", "text",
                                 "html", "image", or "unknown"
                text: str | None — extracted text (None for images)
                image_b64: str | None — base64-encoded image bytes, None otherwise
                media_type: str | None — e.g. "image/png" for images, None otherwise
                size_info: str — human-readable size ("12 pages", "847 rows", "2.3 KB")
                truncated: bool — whether content was cut off
                truncation_notice: str | None — explanation of truncation
                needs_storage_prompt: bool — whether to show storage menu
        """
        file_type = self._detect_type(filename, content_type)

        try:
            if file_type == "pdf":
                return await self._extract_pdf(filename, content_bytes)
            elif file_type == "word":
                return await self._extract_word(filename, content_bytes)
            elif file_type == "powerpoint":
                return await self._extract_powerpoint(filename, content_bytes)
            elif file_type == "csv":
                return await self._extract_csv(filename, content_bytes)
            elif file_type == "text":
                return await self._extract_text(filename, content_bytes)
            elif file_type == "html":
                return await self._extract_html(filename, content_bytes)
            elif file_type == "image":
                return await self._extract_image(filename, content_bytes, content_type)
            else:
                ext = os.path.splitext(filename)[1].lower() if filename else ""
                return self._build_result(
                    file_name=filename,
                    file_type="unknown",
                    text=(
                        f"Unsupported file type: {ext}. "
                        "Supported: PDF, Word (.docx), PowerPoint (.pptx), CSV, TXT, MD, HTML, "
                        "and images (PNG, JPG, JPEG, WEBP, GIF)."
                    ),
                    image_b64=None,
                    media_type=None,
                    size_info="unknown",
                    truncated=False,
                    truncation_notice=None,
                    needs_storage_prompt=False,
                )
        except Exception as e:
            logger.error(f"DocumentExtractionService.extract({filename!r}): {repr(e)}")
            return self._build_result(
                file_name=filename,
                file_type=file_type,
                text=f"[Extraction failed: {repr(e)}]",
                image_b64=None,
                media_type=None,
                size_info="unknown",
                truncated=False,
                truncation_notice=None,
                needs_storage_prompt=False,
            )

    # -------------------------------------------------------------------------
    # Type detection
    # -------------------------------------------------------------------------

    def _detect_type(self, filename: str, content_type: str) -> str:
        """
        Determine the logical file type.
        Extension takes priority over content_type.
        """
        if filename:
            ext = os.path.splitext(filename)[1].lower()
            if ext in self.EXTENSION_MAP:
                return self.EXTENSION_MAP[ext]
            if ext in self.IMAGE_EXTENSIONS:
                return "image"

        # Fallback: content_type
        if content_type:
            ct = content_type.lower()
            if "pdf" in ct:
                return "pdf"
            if "word" in ct or "docx" in ct or "document" in ct:
                return "word"
            if "powerpoint" in ct or "pptx" in ct or "presentation" in ct:
                return "powerpoint"
            if "csv" in ct:
                return "csv"
            if ct.startswith("text/"):
                return "text"
            if "html" in ct:
                return "html"
            if ct.startswith("image/"):
                return "image"

        return "unknown"

    # -------------------------------------------------------------------------
    # PDF extraction
    # -------------------------------------------------------------------------

    async def _extract_pdf(self, filename: str, content_bytes: bytes) -> dict:
        try:
            import pdfplumber
        except ImportError as e:
            return self._build_result(
                file_name=filename, file_type="pdf",
                text=f"[Extraction failed: {repr(e)}]",
                image_b64=None, media_type=None,
                size_info="unknown", truncated=False,
                truncation_notice=None, needs_storage_prompt=False,
            )

        try:
            with pdfplumber.open(io.BytesIO(content_bytes)) as pdf:
                total_pages = len(pdf.pages)
                truncated = total_pages > 20
                pages_to_read = pdf.pages[:20] if truncated else pdf.pages

                parts = []
                for page in pages_to_read:
                    page_text = page.extract_text()
                    if page_text:
                        parts.append(page_text)

                text = "\n\n".join(parts)

                size_info = "1 page" if total_pages == 1 else f"{total_pages} pages"
                truncation_notice = None
                if truncated:
                    truncation_notice = (
                        f"This PDF is {total_pages} pages — I've read the first 20 pages. "
                        "Let me know if you'd like me to focus on a specific section."
                    )

                needs_storage_prompt = truncated or total_pages > 5

                return self._build_result(
                    file_name=filename, file_type="pdf",
                    text=text, image_b64=None, media_type=None,
                    size_info=size_info, truncated=truncated,
                    truncation_notice=truncation_notice,
                    needs_storage_prompt=needs_storage_prompt,
                )
        except Exception as e:
            logger.error(f"PDF extraction failed for {filename!r}: {repr(e)}")
            return self._build_result(
                file_name=filename, file_type="pdf",
                text=f"[Extraction failed: {repr(e)}]",
                image_b64=None, media_type=None,
                size_info="unknown", truncated=False,
                truncation_notice=None, needs_storage_prompt=False,
            )

    # -------------------------------------------------------------------------
    # Word (.docx) extraction
    # -------------------------------------------------------------------------

    async def _extract_word(self, filename: str, content_bytes: bytes) -> dict:
        try:
            from docx import Document
        except ImportError as e:
            return self._build_result(
                file_name=filename, file_type="word",
                text=f"[Extraction failed: {repr(e)}]",
                image_b64=None, media_type=None,
                size_info="unknown", truncated=False,
                truncation_notice=None, needs_storage_prompt=False,
            )

        try:
            doc = Document(io.BytesIO(content_bytes))
            paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
            full_text = "\n".join(paragraphs)

            # Estimate pages: every 3000 chars ≈ 1 page
            total_chars = len(full_text)
            approx_pages = max(1, (total_chars + 2999) // 3000)

            truncated = False
            truncation_notice = None
            max_chars = 20 * 3000  # 60,000 chars

            if total_chars > max_chars:
                full_text = full_text[:max_chars]
                truncated = True
                truncation_notice = (
                    f"This document is approximately {approx_pages} pages — "
                    "I've read the first 20 pages. "
                    "Let me know if you'd like me to focus on a specific section."
                )

            displayed_pages = max(1, (len(full_text) + 2999) // 3000)
            size_info = "1 page" if approx_pages == 1 else f"{approx_pages} pages"
            needs_storage_prompt = truncated or approx_pages > 5

            return self._build_result(
                file_name=filename, file_type="word",
                text=full_text, image_b64=None, media_type=None,
                size_info=size_info, truncated=truncated,
                truncation_notice=truncation_notice,
                needs_storage_prompt=needs_storage_prompt,
            )
        except Exception as e:
            logger.error(f"Word extraction failed for {filename!r}: {repr(e)}")
            return self._build_result(
                file_name=filename, file_type="word",
                text=f"[Extraction failed: {repr(e)}]",
                image_b64=None, media_type=None,
                size_info="unknown", truncated=False,
                truncation_notice=None, needs_storage_prompt=False,
            )

    # -------------------------------------------------------------------------
    # PowerPoint (.pptx / .ppt) extraction
    # -------------------------------------------------------------------------

    async def _extract_powerpoint(self, filename: str, content_bytes: bytes) -> dict:
        try:
            from pptx import Presentation
        except ImportError as e:
            return self._build_result(
                file_name=filename, file_type="powerpoint",
                text=f"[Extraction failed: {repr(e)}]",
                image_b64=None, media_type=None,
                size_info="unknown", truncated=False,
                truncation_notice=None, needs_storage_prompt=False,
            )

        try:
            prs = Presentation(io.BytesIO(content_bytes))
            total_slides = len(prs.slides)
            truncated = total_slides > 20
            slides_to_read = list(prs.slides)[:20] if truncated else prs.slides

            parts = []
            for i, slide in enumerate(slides_to_read, start=1):
                slide_texts = []
                for shape in slide.shapes:
                    if hasattr(shape, "text_frame"):
                        for para in shape.text_frame.paragraphs:
                            line = para.text.strip()
                            if line:
                                slide_texts.append(line)
                if slide_texts:
                    parts.append(f"[Slide {i}]\n" + "\n".join(slide_texts))

            text = "\n\n".join(parts)
            size_info = "1 slide" if total_slides == 1 else f"{total_slides} slides"

            truncation_notice = None
            if truncated:
                truncation_notice = (
                    f"This presentation has {total_slides} slides — I've read the first 20. "
                    "Let me know if you'd like me to focus on a specific section."
                )

            needs_storage_prompt = truncated or total_slides > 5

            return self._build_result(
                file_name=filename, file_type="powerpoint",
                text=text, image_b64=None, media_type=None,
                size_info=size_info, truncated=truncated,
                truncation_notice=truncation_notice,
                needs_storage_prompt=needs_storage_prompt,
            )
        except Exception as e:
            logger.error(f"PowerPoint extraction failed for {filename!r}: {repr(e)}")
            return self._build_result(
                file_name=filename, file_type="powerpoint",
                text=f"[Extraction failed: {repr(e)}]",
                image_b64=None, media_type=None,
                size_info="unknown", truncated=False,
                truncation_notice=None, needs_storage_prompt=False,
            )

    # -------------------------------------------------------------------------
    # CSV extraction
    # -------------------------------------------------------------------------

    async def _extract_csv(self, filename: str, content_bytes: bytes) -> dict:
        try:
            text_data = content_bytes.decode("utf-8", errors="replace")
            reader = csv.reader(io.StringIO(text_data))
            rows = list(reader)

            total_rows = len(rows)
            truncated = total_rows > 500
            rows_to_use = rows[:500] if truncated else rows

            # Format as pipe-separated table
            lines = []
            for i, row in enumerate(rows_to_use):
                lines.append(" | ".join(str(cell) for cell in row))

            text = "\n".join(lines)
            size_info = "1 row" if total_rows == 1 else f"{total_rows} rows"

            truncation_notice = None
            if truncated:
                truncation_notice = (
                    f"This CSV has {total_rows} rows — I've read the first 500."
                )

            # Large CSV or truncated = ask about storage
            needs_storage_prompt = truncated or len(content_bytes) > 50 * 1024

            return self._build_result(
                file_name=filename, file_type="csv",
                text=text, image_b64=None, media_type=None,
                size_info=size_info, truncated=truncated,
                truncation_notice=truncation_notice,
                needs_storage_prompt=needs_storage_prompt,
            )
        except Exception as e:
            logger.error(f"CSV extraction failed for {filename!r}: {repr(e)}")
            return self._build_result(
                file_name=filename, file_type="csv",
                text=f"[Extraction failed: {repr(e)}]",
                image_b64=None, media_type=None,
                size_info="unknown", truncated=False,
                truncation_notice=None, needs_storage_prompt=False,
            )

    # -------------------------------------------------------------------------
    # Plain text / Markdown extraction
    # -------------------------------------------------------------------------

    async def _extract_text(self, filename: str, content_bytes: bytes) -> dict:
        try:
            text = content_bytes.decode("utf-8", errors="replace")
            total_chars = len(text)
            truncated = total_chars > 100_000

            if truncated:
                text = text[:100_000]

            approx_pages = max(1, (total_chars + 2999) // 3000)
            if total_chars < 3000:
                size_info = f"{total_chars} characters"
            else:
                size_info = f"{approx_pages} pages"

            truncation_notice = None
            if truncated:
                truncation_notice = (
                    f"This file is {total_chars} characters — I've read the first 100,000. "
                    "Let me know if you'd like me to focus on a specific section."
                )

            needs_storage_prompt = truncated or total_chars > 50_000

            return self._build_result(
                file_name=filename, file_type="text",
                text=text, image_b64=None, media_type=None,
                size_info=size_info, truncated=truncated,
                truncation_notice=truncation_notice,
                needs_storage_prompt=needs_storage_prompt,
            )
        except Exception as e:
            logger.error(f"Text extraction failed for {filename!r}: {repr(e)}")
            return self._build_result(
                file_name=filename, file_type="text",
                text=f"[Extraction failed: {repr(e)}]",
                image_b64=None, media_type=None,
                size_info="unknown", truncated=False,
                truncation_notice=None, needs_storage_prompt=False,
            )

    # -------------------------------------------------------------------------
    # HTML extraction
    # -------------------------------------------------------------------------

    async def _extract_html(self, filename: str, content_bytes: bytes) -> dict:
        try:
            html = content_bytes.decode("utf-8", errors="replace")
            # Strip HTML tags with a simple regex
            text = re.sub(r"<[^>]+>", "", html)
            # Collapse excessive whitespace
            text = re.sub(r"\n{3,}", "\n\n", text).strip()

            total_chars = len(text)
            truncated = total_chars > 100_000
            if truncated:
                text = text[:100_000]

            size_info = f"{total_chars} characters"

            truncation_notice = None
            if truncated:
                truncation_notice = (
                    f"This HTML file is {total_chars} characters — I've read the first 100,000."
                )

            needs_storage_prompt = truncated or total_chars > 50_000

            return self._build_result(
                file_name=filename, file_type="html",
                text=text, image_b64=None, media_type=None,
                size_info=size_info, truncated=truncated,
                truncation_notice=truncation_notice,
                needs_storage_prompt=needs_storage_prompt,
            )
        except Exception as e:
            logger.error(f"HTML extraction failed for {filename!r}: {repr(e)}")
            return self._build_result(
                file_name=filename, file_type="html",
                text=f"[Extraction failed: {repr(e)}]",
                image_b64=None, media_type=None,
                size_info="unknown", truncated=False,
                truncation_notice=None, needs_storage_prompt=False,
            )

    # -------------------------------------------------------------------------
    # Image handling (multimodal passthrough)
    # -------------------------------------------------------------------------

    async def _extract_image(self, filename: str, content_bytes: bytes, content_type: str) -> dict:
        try:
            image_b64 = base64.b64encode(content_bytes).decode("ascii")

            # Determine media_type from content_type, fallback to extension
            media_type = content_type if content_type and content_type.startswith("image/") else None
            if not media_type and filename:
                ext = os.path.splitext(filename)[1].lower()
                ext_to_mime = {
                    ".png": "image/png",
                    ".jpg": "image/jpeg",
                    ".jpeg": "image/jpeg",
                    ".webp": "image/webp",
                    ".gif": "image/gif",
                }
                media_type = ext_to_mime.get(ext)

            size_kb = len(content_bytes) / 1024
            if size_kb < 1024:
                size_info = f"{size_kb:.1f} KB"
            else:
                size_info = f"{size_kb / 1024:.1f} MB"

            return self._build_result(
                file_name=filename, file_type="image",
                text=None, image_b64=image_b64, media_type=media_type,
                size_info=size_info, truncated=False,
                truncation_notice=None,
                needs_storage_prompt=False,  # Images always ephemeral by default
            )
        except Exception as e:
            logger.error(f"Image handling failed for {filename!r}: {repr(e)}")
            return self._build_result(
                file_name=filename, file_type="image",
                text=f"[Extraction failed: {repr(e)}]",
                image_b64=None, media_type=None,
                size_info="unknown", truncated=False,
                truncation_notice=None, needs_storage_prompt=False,
            )

    # -------------------------------------------------------------------------
    # Helper
    # -------------------------------------------------------------------------

    def _build_result(
        self,
        file_name: str,
        file_type: str,
        text,
        image_b64,
        media_type,
        size_info: str,
        truncated: bool,
        truncation_notice,
        needs_storage_prompt: bool,
    ) -> dict:
        return {
            "file_name": file_name,
            "file_type": file_type,
            "text": text,
            "image_b64": image_b64,
            "media_type": media_type,
            "size_info": size_info,
            "truncated": truncated,
            "truncation_notice": truncation_notice,
            "needs_storage_prompt": needs_storage_prompt,
        }
